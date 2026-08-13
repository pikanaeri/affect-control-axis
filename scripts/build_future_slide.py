#!/usr/bin/env python
"""Single slide: future directions + methods overview + linked LLM/VLM precedents."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_REPO = pathlib.Path(__file__).resolve().parents[1]
OUT = str(_REPO / "slides" / "future_directions_slide.pptx")

NAVY=RGBColor(0x1E,0x27,0x61); DKNAVY=RGBColor(0x12,0x16,0x33); ICE=RGBColor(0xCA,0xDC,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); BLUE=RGBColor(0x00,0x72,0xB2); ORANGE=RGBColor(0xE6,0x9F,0x00)
GREEN=RGBColor(0x00,0x9E,0x73); INK=RGBColor(0x1A,0x1A,0x2E); MUTED=RGBColor(0x6B,0x72,0x80)
CARD=RGBColor(0xF2,0xF5,0xFB); LINK=RGBColor(0x0B,0x5C,0xAD)
HEAD="Cambria"; BODY="Calibri"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6])
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb=WHITE; bg.line.fill.background(); bg.shadow.inherit=False
s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2,bg._element)

def card(x,y,w,h,fill=CARD):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); sh.shadow.inherit=False
    try: sh.adjustments[0]=0.03
    except Exception: pass
    return sh

def pill(x,y,text,fill=BLUE,tc=WHITE,sz=11):
    w=0.11*len(text)+0.3; sh=card(x,y,w,0.32,fill=fill)
    try: sh.adjustments[0]=0.5
    except Exception: pass
    tf=sh.text_frame; tf.word_wrap=False
    for m in ("left","right","top","bottom"): setattr(tf,f"margin_{m}",Inches(0.08) if m in("left","right") else 0)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=True; r.font.name=BODY; r.font.color.rgb=tc

def txt(x,y,w,h,paras,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ("left","right","top","bottom"): setattr(tf,f"margin_{m}",0)
    for i,p in enumerate(paras):
        para=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.space_after=Pt(p.get("sa",3)); para.line_spacing=p.get("ls",1.02)
        for run in p["runs"]:
            r=para.add_run(); r.text=run["t"]
            f=r.font; f.size=Pt(run.get("sz",11)); f.bold=run.get("b",False); f.italic=run.get("i",False)
            f.name=run.get("font",BODY); f.color.rgb=run.get("c",INK)
            if run.get("link"): r.hyperlink.address=run["link"]; f.color.rgb=LINK
    return tb

# ---- header ----
pill(0.6,0.42,"FUTURE DIRECTIONS · METHODS OVERVIEW",fill=NAVY,sz=11)
txt(0.58,0.85,12.2,0.7,[{"runs":[{"t":"Affect-specific behavior probes — does the VLM track emotion appraisals, not just valence?","sz":24,"b":True,"font":HEAD,"c":NAVY}],"ls":1.0}])

# ---- methods overview strip ----
card(0.6,1.72,12.13,0.78,fill=CARD)
txt(0.85,1.82,11.7,0.62,[
 {"runs":[{"t":"Method (every probe): ","b":True,"sz":12,"c":INK},
          {"t":"emotion vector via diff-in-means → steer / text-prime → first-token option-logit → causal mediation.","sz":12,"c":INK}],"sa":2},
 {"runs":[{"t":"Libraries (code-backed): ","b":True,"sz":12,"c":INK},
          {"t":"TransformerLens · pyvene · steering-vectors · scikit-learn.   Readout = log-odds of answer A vs B, generation-free.","sz":12,"c":INK}]}])

# ---- left column: future probes ----
LX,LW=0.6,6.05
card(LX,2.68,LW,4.02,fill=CARD)
pill(LX+0.25,2.9,"FUTURE PROBES  (appraisal specificity)",fill=GREEN,sz=11)
def grp(y,title,items):
    txt(LX+0.28,y,LW-0.5,0.3,[{"runs":[{"t":title,"b":True,"sz":11.5,"c":NAVY}]}])
    yy=y+0.32
    for name,cite in items:
        txt(LX+0.34,yy,LW-0.6,0.3,[{"runs":[{"t":"• "+name+"  ","sz":11,"c":INK},{"t":"["+cite+"]","i":True,"sz":9.5,"c":MUTED}]}])
        yy+=0.262
    return yy
y=grp(3.42,"Specificity — same valence, different behavior",[
 ("Fear vs anger → risk (opposite)","Lerner & Keltner 2001"),
 ("Sad vs anxious → risk / reward","Raghunathan & Pham 1999"),
 ("Sad vs anger → attribution","Keltner et al. 1993")])
y=grp(y+0.05,"Emotion-specific decisions",[
 ("Sadness → impatience","Lerner, Li & Weber 2013"),
 ("Sadness → valuation (buy/sell)","Lerner et al. 2004"),
 ("Anger → punitiveness","Lerner et al. 1998"),
 ("Affect → generalized risk","Johnson & Tversky 1983")])
y=grp(y+0.05,"Moral / social emotions",[
 ("Trust · Fairness / ultimatum","Dunn 2005 · Harlé 2007"),
 ("Gratitude · Guilt","DeSteno 2010 · Ketelaar 2003")])

# ---- right column: LLM/VLM precedents (linked) ----
RX,RW=6.9,5.83
card(RX,2.68,RW,4.02,fill=CARD)
pill(RX+0.25,2.9,"PRIOR LLM / VLM BEHAVIOR ANALYSES  (method precedents)",fill=BLUE,sz=11)
def prec(y,author,desc,url):
    txt(RX+0.3,y,RW-0.55,0.62,[
      {"runs":[{"t":author+" — ","b":True,"sz":11,"c":INK},{"t":desc+"  ","sz":11,"c":INK},
               {"t":url.split("//")[-1],"sz":10,"link":url}]}],)
    return y+0.72
y=3.42
y=prec(y,"Coda-Forno et al. 2023","inducing anxiety shifts LLM exploration & bias (affect→behavior, code)","https://arxiv.org/abs/2304.11111")
y=prec(y,"Binz & Schulz 2023 (PNAS)","cognitive psychology of GPT-3: decisions, risk, biases","https://arxiv.org/abs/2206.14576")
y=prec(y,"Aher et al. 2023 (ICML)","LLMs replicate human studies incl. Ultimatum Game (code)","https://arxiv.org/abs/2208.10264")
y=prec(y,"Arditi et al. 2024","steer / ablate a behavior direction — our core method","https://arxiv.org/abs/2406.11717")
txt(RX+0.3,y,RW-0.55,0.5,[{"runs":[
   {"t":"Anthropic 2026 — ","b":True,"sz":11,"c":INK},
   {"t":"emotion vectors → behavior (text-only; we extend to images).","sz":11,"c":INK}]}])

# ---- footer thesis ----
txt(0.6,6.9,12.2,0.55,[{"runs":[
 {"t":"Goal: ","b":True,"sz":11.5,"c":ORANGE},
 {"t":"show the VLM tracks emotion appraisals (certainty · control · agency), not just valence — extending affect→behavior LLM work to images + internal mediation.","sz":11.5,"c":INK}],"ls":1.0}])

prs.save(OUT); print("saved:",OUT)
