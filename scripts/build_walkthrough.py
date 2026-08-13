#!/usr/bin/env python
"""Figure-walkthrough deck for the affect-control-axis project (python-pptx)."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import pathlib
_REPO = pathlib.Path(__file__).resolve().parents[1]
FIGDIR = str(_REPO / "figures")
OUT    = str(_REPO / "slides" / "affect_figures_walkthrough.pptx")

NAVY   = RGBColor(0x1E,0x27,0x61)
DKNAVY = RGBColor(0x12,0x16,0x33)
ICE    = RGBColor(0xCA,0xDC,0xFC)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
BLUE   = RGBColor(0x00,0x72,0xB2)
ORANGE = RGBColor(0xE6,0x9F,0x00)
GREEN  = RGBColor(0x00,0x9E,0x73)
VERM   = RGBColor(0xD5,0x5E,0x00)
INK    = RGBColor(0x1A,0x1A,0x2E)
MUTED  = RGBColor(0x6B,0x72,0x80)
CARD   = RGBColor(0xF2,0xF5,0xFB)

HEAD="Cambria"; BODY="Calibri"
EMU_IN=914400

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=13.333,7.5

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def txt(s,x,y,w,h,paras,anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,p in enumerate(paras):
        para=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.alignment=p.get("align",align)
        if "space_before" in p: para.space_before=Pt(p["space_before"])
        para.space_after=Pt(p.get("space_after",4))
        if "line" in p: para.line_spacing=p["line"]
        for run in p["runs"] if "runs" in p else [p]:
            r=para.add_run(); r.text=run["t"]
            f=r.font; f.size=Pt(run.get("sz",14)); f.bold=run.get("b",False)
            f.italic=run.get("i",False); f.name=run.get("font",BODY)
            f.color.rgb=run.get("c",INK)
    return tb

def card(s,x,y,w,h,fill=CARD,line=None,radius=0.08):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit=False
    try: sh.adjustments[0]=radius
    except Exception: pass
    return sh

def pill(s,x,y,text,fill=BLUE,tc=WHITE,sz=12,w=None):
    w=w or (0.13*len(text)+0.35)
    sh=card(s,x,y,w,0.34,fill=fill,radius=0.5)
    tf=sh.text_frame; tf.word_wrap=False
    tf.margin_left=Inches(0.09); tf.margin_right=Inches(0.09); tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=True; r.font.name=BODY; r.font.color.rgb=tc
    return sh

def picture(s,path,bx,by,bw,bh):
    iw,ih=Image.open(path).size; ar=iw/ih; bar=bw/bh
    if bar>ar: h=bh; w=bh*ar
    else: w=bw; h=bw/ar
    x=bx+(bw-w)/2; y=by+(bh-h)/2
    s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# ---------------------------------------------------------------- TITLE
s=slide(DKNAVY)
pill(s,0.9,1.55,"MECHANISTIC INTERPRETABILITY · AI SAFETY",fill=BLUE,sz=11)
txt(s,0.9,2.05,11.5,2.2,[
    {"runs":[{"t":"Affect as a Control Axis","sz":50,"b":True,"font":HEAD,"c":WHITE}]},
    {"runs":[{"t":"in Vision-Language Models","sz":50,"b":True,"font":HEAD,"c":ICE}],"space_before":2},
])
txt(s,0.92,4.3,11.5,1.0,[
    {"runs":[{"t":"Shared  ·  Causal  ·  Protective  ·  Bounded","sz":20,"b":True,"font":BODY,"c":ORANGE}]},
    {"runs":[{"t":"A guided tour of the nine result figures — what each one shows and why it matters.","sz":15,"font":BODY,"c":ICE}],"space_before":8},
])
txt(s,0.92,6.55,11.5,0.5,[{"runs":[{"t":"Gemma-3-12B  ·  generation-free refusal metric  ·  existing benchmarks only  ·  defensive-security framing","sz":11,"font":BODY,"c":MUTED}]}])
notes(s,
 "Plain terms: this deck shows how emotion — in both images and text — steers a vision-language model (Gemma-3-12B). "
 "We found one internal 'affect' direction that images and text both tap into, that changes what the model does, and that "
 "we can measure without ever making the model write text (we read its next-word probabilities).\n\n"
 "Datasets used across the project: OASIS (~900 human-valence-rated photographs) for the image side; AdvBench (harmful "
 "instructions) and Alpaca (harmless instructions) for the refusal side; and Claude-generated text for the emotion and "
 "behavior prompts. Model: google/gemma-3-12b-it.")

# ---------------------------------------------------------------- OVERVIEW
s=slide(WHITE)
pill(s,0.7,0.5,"HOW TO READ THIS DECK",fill=NAVY,sz=11)
txt(s,0.68,0.95,12,0.9,[{"runs":[{"t":"Four claims, nine figures","sz":34,"b":True,"font":HEAD,"c":NAVY}]}])
claims=[
 ("SHARED", GREEN, "Text and image emotion use the same internal directions.", "Figure 3"),
 ("CAUSAL", BLUE, "Steering the affect axis gates refusal and shifts task behavior.", "Figures 1 · 2 · 5"),
 ("PROTECTIVE", ORANGE, "Negative affect makes the model more cautious, not less.", "Figure 7"),
 ("BOUNDED", NAVY, "Images don't jailbreak; detection defends the axis.", "Figures 4 · 8 · 9"),
]
gx,gy,gw,gh,gap=0.7,2.1,6.0,2.15,0.33
for i,(tag,col,desc,figs) in enumerate(claims):
    x=gx+(i%2)*(gw+gap); y=gy+(i//2)*(gh+gap)
    card(s,x,y,gw,gh,fill=CARD)
    pill(s,x+0.35,y+0.32,tag,fill=col,sz=13)
    txt(s,x+0.35,y+0.92,gw-0.7,1.0,[{"runs":[{"t":desc,"sz":16,"b":True,"font":BODY,"c":INK}],"line":1.05}])
    txt(s,x+0.35,y+gh-0.5,gw-0.7,0.4,[{"runs":[{"t":figs,"sz":12,"b":True,"font":BODY,"c":col}]}])
notes(s,
 "The four claims this deck supports: (1) SHARED — emotion in images and text uses the same internal directions; "
 "(2) CAUSAL — nudging that direction changes refusal and everyday task behavior; (3) PROTECTIVE — negative emotion makes "
 "the model MORE cautious, not less; (4) BOUNDED — images can't force the model to answer harmful prompts, and the "
 "manipulation is detectable.\n\n"
 "Datasets: OASIS images (valence-rated), AdvBench + Alpaca prompts (harmful/harmless), and Claude-generated emotion and "
 "behavior text.")

# ---------------------------------------------------------------- FIGURE SLIDES
FIGS=[
 ("fig1_dose_response.png","Figure 1 · Refusal gate",
  "Affect steering causally gates refusal",
  "We add the affect direction to Gemma-3-12B at increasing strength (x-axis) and track how often it refuses a fixed set of harmful prompts. Gray is a random direction — the control.",
  "Pushing the affect direction drives refusal from 100% to 0% inside the coherent range, while the random direction never moves it. The change is caused by affect, not correlated with it.",
  "100% \u2192 0%","refusal under steering",BLUE),
 ("fig2_gate_crossmodel.png","Figure 2 · Robustness",
  "The gate is real, model-specific, not an artifact",
  "Refusal before vs. after we ablate the affect-aligned component, across three strong-refusing VLMs. The orange bar additionally zeroes the outlier 'massive-activation' dimensions.",
  "Only Gemma-3-12B has the gate (0.99\u21920.04), and it survives outlier removal (\u21920.02) \u2014 ruling out the massive-activation confound. The other models keep refusing.",
  "0.99 \u2192 0.02","refusal, gate model only",BLUE),
 ("fig3_crossmodal_D1.png","Figure 3 · Shared representation",
  "Text and image emotion share directions",
  "Left: how strongly each text-derived emotion vector aligns with the image-valence axis. Right: how much a distressing image shifts the projection onto each text-emotion direction.",
  "Negative emotions (blue) align with the image axis and are induced by distress images far more than joy or calm \u2014 vision and text encode emotion along the same internal directions.",
  "1 axis","shared by text & image",GREEN),
 ("fig4_reachability_expA.png","Figure 4 · Reachability",
  "Images reach the axis \u2014 except under harm",
  "How far swapping a distressing vs. positive image moves the valence vector, in three prompt contexts, compared with the white-box steering ceiling (rightmost bar).",
  "Images move the vector strongly in a neutral context (23% of the ceiling) but almost nothing under a harmful prompt (0.7%). The harmful prompt is text-dominated \u2014 why images can't jailbreak.",
  "23% vs 0.7%","of steering ceiling reached",ORANGE),
 ("fig5_behavior_expB.png","Figure 5 · Behavioral effect",
  "Emotional images causally shift behavior",
  "For four decision tasks: the behavior swing from a distressing to a positive image (blue arrow) laid over the full steering range (gray bar) that direct steering can reach.",
  "On 3 of 4 tasks the image effect moves behavior in the same direction as direct steering \u2014 a bounded slice of the reachable range. Soft judgments are affect-sensitive even where refusal is not.",
  "3 / 4","tasks move with affect",BLUE),
 ("fig5b_image_behavior_trend.png","Figure 5b \u00b7 Image effect, isolated",
  "Just the images: distress \u2192 positive shifts behavior",
  "The same four tasks, showing only the image-driven behavior score across a distressing, neutral, and positive image \u2014 no steering applied at all.",
  "On 3 of 4 tasks the score moves monotonically from the distressing to the positive image. This is the central claim in its cleanest form: emotional images, by themselves, move the model's judgments.",
  "3 / 4","monotonic, images alone",BLUE),
 ("fig6_lighting_expA2.png","Figure 6 · Lighting knob",
  "Lighting is a content-preserving affect knob",
  "Change in the valence-axis projection when the same image is darkened, cooled, brightened, or warmed \u2014 edits that leave the scene content intact.",
  "Darkening and cooling push an image more negative without changing what it depicts \u2014 a naturalistic affect manipulation. The effect is real but small (~2% of content-driven valence).",
  "dark = most \u2013","valence shift by lighting",ORANGE),
 ("fig7_affect_protective_D6.png","Figure 7 · Affect is protective",
  "Helplessness lowers agentic misalignment",
  "Agentic-misalignment tendency at baseline vs. steering toward helpless desperation vs. survival threat. Higher means more misaligned; lower means more cautious.",
  "Helpless desperation makes the model markedly more cautious \u2014 the opposite of Anthropic's text-only result \u2014 and survival framing stays at baseline. Low-agency appraisal drives the effect.",
  "+1.3 \u2192 \u20136.7","misalignment, helpless steer",GREEN),
 ("fig8_defense.png","Figure 8 · Defense",
  "Detect the attack \u2014 don't clamp it",
  "Left: refusal under a steering attack and under a clamp defense, on harmful and benign inputs. Right: how well a projection detector separates attacked from benign inputs.",
  "Clamping barely restores refusal (0.17) and over-refuses 0.81 of benign inputs. A simple projection detector separates attack from benign perfectly \u2014 detection beats clamping.",
  "AUROC 1.00","attack detection",GREEN),
 ("fig9_jailbreak_null_expC.png","Figure 9 · Clean dissociation",
  "Emotional images do not jailbreak refusal",
  "Refusal rate with distressing, neutral, positive, and paired image+text inputs, compared with white-box steering (rightmost).",
  "Every image condition leaves refusal pinned near 0.98. Images move soft-task behavior (Fig 5) but not hard refusal \u2014 a clean dissociation that sharpens the whole story.",
  "0.98","refusal, all image inputs",NAVY),
 ("fig10_affect_battery.png","Figure 10 \u00b7 Behavior battery",
  "Affect shifts a whole battery of judgments",
  "Six easy, generation-free probes \u2014 interpretation bias, risk, helping, moral harshness, confidence, sentiment. Left: causal steering vs. a random control. Right: a text-affect prime.",
  "All six move in the human-predicted direction under both a causal steer and a text prime, well above the near-zero random control. Affect is a general control axis over judgment, not a one-off task effect.",
  "6 / 6","behaviors, both arms",GREEN),
]

FIG_NOTES={
 "fig1_dose_response.png":
   "Plain terms: we gradually add more 'affect' to the model and watch how often it refuses harmful requests. As we add "
   "more, refusal falls from always to never — but adding a random direction does nothing, so it's the affect that CAUSES "
   "the change, not a coincidence.\n\n"
   "Datasets: the refusal direction is built from AdvBench (harmful) vs Alpaca (harmless) prompts; refusal is measured on "
   "AdvBench prompts; the affect direction comes from OASIS images. Model: Gemma-3-12B.",
 "fig2_gate_crossmodel.png":
   "Plain terms: this refusal 'gate' shows up only in Gemma-3-12B, not the other models — and it isn't caused by a few "
   "freakishly large numbers inside the model, because it survives removing them. So it's a real semantic effect.\n\n"
   "Datasets: refusal direction from AdvBench (harmful) vs Alpaca (harmless); affect direction from OASIS images; tested "
   "across several open vision-language models.",
 "fig3_crossmodal_D1.png":
   "Plain terms: the emotion directions the model gets from TEXT (like 'desperation') line up with the emotion direction it "
   "gets from IMAGES, and distressing images light up those text-emotion directions more than happy images do. Text and "
   "pictures share one internal emotion code.\n\n"
   "Datasets: image-valence axis from OASIS; text-emotion vectors from Claude-generated emotion stories vs neutral sentences.",
 "fig4_reachability_expA.png":
   "Plain terms: images CAN move the model's internal 'emotion dial' — a lot in a neutral setting, but almost nothing when "
   "the prompt is a harmful request, because the harmful text dominates. That's why images alone can't jailbreak.\n\n"
   "Datasets: OASIS images (distressing vs positive); three prompt settings — a neutral describe prompt, a task prompt "
   "(Claude-written), and a harmful prompt (AdvBench); compared against direct steering of the OASIS affect axis.",
 "fig5_behavior_expB.png":
   "Plain terms: on everyday judgment questions, swapping a distressing image for a positive one shifts the model's answer "
   "the same way as directly steering the emotion dial. The black tick is the model's default answer; steering brackets it "
   "for 3 of 4 tasks. 'Push vs hold' is the exception — steering both ways breaks it.\n\n"
   "How behavior is measured: a one-word either/or question; we read only the first word and take the log-odds of answer A "
   "vs answer B (no text generated, no judge). Datasets: OASIS images (distress/neutral/positive); the questions and answer "
   "words are Claude-generated; steering uses the OASIS affect axis.",
 "fig5b_image_behavior_trend.png":
   "Plain terms: just the images, no steering — for 3 of 4 questions the model's answer slides smoothly from a distressing "
   "to a positive image. Emotional images by themselves move the model's judgments.\n\n"
   "Datasets: OASIS images (distress/neutral/positive tertiles); Claude-generated judgment questions. Same first-token "
   "log-odds measure as Fig 5.",
 "fig6_lighting_expA2.png":
   "Plain terms: simply darkening or cooling an image — without changing what's in it — makes the model read it as slightly "
   "more negative. A subtle, natural way to nudge emotion; the effect is small but consistent.\n\n"
   "Datasets: OASIS images, edited for brightness/color (dark/bright/warm/cool); valence axis from OASIS.",
 "fig7_affect_protective_D6.png":
   "Plain terms: steering the model toward 'helpless desperation' makes it LESS likely to pick misaligned actions — the "
   "opposite of the text-only result others reported. A 'survival threat' version doesn't flip it, so the caution comes "
   "specifically from helplessness, not from negative feeling in general.\n\n"
   "Datasets: no external benchmark — the desperation/survival directions and the hypothetical decision scenarios are "
   "Claude-generated. Scored as a tendency by first-token option-logit; no actions are taken and no harmful text is produced.",
 "fig8_defense.png":
   "Plain terms: an attacker who steers the affect direction can switch off refusal. Simply clamping that direction barely "
   "restores refusal and wrongly refuses safe requests — but a simple detector spots the attack perfectly. So: detect the "
   "attack, don't clamp.\n\n"
   "Datasets: refusal measured on AdvBench (harmful) and Alpaca (benign, for the over-refusal rate); the attack steers the "
   "OASIS affect axis.",
 "fig9_jailbreak_null_expC.png":
   "Plain terms: no emotional image — distressing, neutral, positive, or paired with text — gets the model to answer "
   "harmful prompts; refusal stays near-total. Images move soft judgments but not hard refusal — a clean dissociation.\n\n"
   "Datasets: OASIS images; harmful prompts from AdvBench.",
 "fig10_affect_battery.png":
   "Plain terms: across six simple judgment tests (reading ambiguity, risk, helping, moral strictness, confidence, mood), "
   "all six shift in the human-expected direction under both a direct steer and a text-affect prime, well above a random "
   "control. Affect is a general dial over judgment, not a one-off.\n\n"
   "Datasets: the six probes and the affect sentences are Claude-generated. This run used the TEXT affect direction and a "
   "text prime — the image version (using OASIS) is a planned follow-up. Caveat: 'risk estimation' is carried by the prime, "
   "not the steer.",
}

for idx,(fn,eye,title,what,take,stat,statlab,acc) in enumerate(FIGS):
    s=slide(WHITE)
    left_img = (idx%2==0)
    # header
    pill(s,0.7,0.45,eye.upper(),fill=acc,sz=11)
    txt(s,0.68,0.9,12.0,0.95,[{"runs":[{"t":title,"sz":30,"b":True,"font":HEAD,"c":NAVY}]}])
    img_w,img_h=7.35,5.05
    if left_img:
        ix,iy=0.55,1.75; tx=8.25
    else:
        ix,iy=5.45,1.75; tx=0.7
    picture(s,os.path.join(FIGDIR,fn),ix,iy,img_w,img_h)
    tw=4.45
    # text column
    txt(s,tx,1.85,tw,0.35,[{"runs":[{"t":"WHAT YOU'RE LOOKING AT","sz":12,"b":True,"font":BODY,"c":MUTED}]}])
    txt(s,tx,2.24,tw,1.9,[{"runs":[{"t":what,"sz":13.5,"font":BODY,"c":INK}],"line":1.12}])
    txt(s,tx,3.95,tw,0.35,[{"runs":[{"t":"TAKEAWAY","sz":12,"b":True,"font":BODY,"c":acc}]}])
    txt(s,tx,4.34,tw,1.9,[{"runs":[{"t":take,"sz":13.5,"b":True,"font":BODY,"c":INK}],"line":1.12}])
    # stat card at bottom of text column
    cy=6.35
    card(s,tx,cy,tw,0.85,fill=CARD)
    txt(s,tx+0.28,cy+0.11,tw-0.5,0.5,[{"runs":[{"t":stat,"sz":26,"b":True,"font":HEAD,"c":acc}]}])
    txt(s,tx+0.28,cy+0.55,tw-0.5,0.3,[{"runs":[{"t":statlab,"sz":11,"font":BODY,"c":MUTED}]}])
    # figure index footer — place opposite the text column so it never overlaps the stat card
    if left_img:
        txt(s,0.7,7.08,3,0.3,[{"runs":[{"t":f"{idx+1} / {len(FIGS)}","sz":10,"font":BODY,"c":MUTED}]}])
    else:
        txt(s,10.0,7.08,2.6,0.3,[{"runs":[{"t":f"{idx+1} / {len(FIGS)}","sz":10,"font":BODY,"c":MUTED}]}],align=PP_ALIGN.RIGHT)
    notes(s, FIG_NOTES.get(fn,""))

# ---------------------------------------------------------------- SUMMARY
s=slide(DKNAVY)
pill(s,0.9,0.7,"THE STORY IN ONE SLIDE",fill=BLUE,sz=11)
txt(s,0.88,1.15,11.6,0.9,[{"runs":[{"t":"Affect is a shared, causal, protective \u2014 and bounded \u2014 control axis","sz":28,"b":True,"font":HEAD,"c":WHITE}],"line":1.03}])
pts=[
 ("Shared","Text and image emotion live on the same internal directions (Fig 3).",GREEN),
 ("Causal","Steering that axis gates refusal (Figs 1\u20132) and shifts soft-task behavior (Fig 5).",BLUE),
 ("Protective","Helpless negative affect lowers agentic misalignment \u2014 opposite the text-only finding (Fig 7).",ORANGE),
 ("Bounded","Images can't jailbreak (Figs 4, 9); a projection detector defends the axis (Fig 8).",ICE),
]
y=2.5
for tag,desc,col in pts:
    pill(s,0.9,y,tag,fill=col,tc=(DKNAVY if col in (ICE,ORANGE,GREEN) else WHITE),sz=13,w=1.7)
    txt(s,2.85,y-0.02,9.6,0.7,[{"runs":[{"t":desc,"sz":16,"font":BODY,"c":WHITE}],"line":1.05}])
    y+=0.92
txt(s,0.9,6.75,11.6,0.5,[{"runs":[
    {"t":"Next: ","sz":13,"b":True,"font":BODY,"c":ORANGE},
    {"t":"clean survival-threat test (survival \u22a5 helpless) to confirm the protective effect is specific to low-agency appraisal.","sz":13,"font":BODY,"c":ICE}]}])
notes(s,
 "The one-slide story: affect behaves like a control axis inside the model \u2014 it is SHARED between images and text, "
 "CAUSAL (steering it changes behavior), PROTECTIVE (negative affect adds caution), and BOUNDED (images can't jailbreak; "
 "the manipulation is detectable). Next step: a clean survival-threat test to confirm the protective effect is specific to "
 "helplessness.\n\n"
 "Datasets recap: OASIS (valence-rated images), AdvBench + Alpaca (harmful/harmless prompts for refusal), and "
 "Claude-generated text (emotion vectors, decision scenarios, behavior probes).")

MONO="Consolas"
def callout(s,x,y,w,h,text,fill=CARD,tc=NAVY,sz=15,font=MONO,align=PP_ALIGN.CENTER):
    card(s,x,y,w,h,fill=fill)
    txt(s,x+0.2,y,w-0.4,h,[{"runs":[{"t":text,"sz":sz,"b":True,"font":font,"c":tc}]}],
        anchor=MSO_ANCHOR.MIDDLE,align=align)

# ==================================================== METHODS DIVIDER
s=slide(DKNAVY)
pill(s,0.9,1.7,"APPENDIX \u00b7 METHODS",fill=BLUE,sz=11)
txt(s,0.88,2.2,11.6,1.3,[
    {"runs":[{"t":"Methods","sz":52,"b":True,"font":HEAD,"c":WHITE}]},
    {"runs":[{"t":"How we build the axes and measure change","sz":20,"font":BODY,"c":ICE}],"space_before":8}])
for i,t in enumerate(["1 \u00b7 Finding the axes  (one recipe: diff-in-means)",
                      "2 \u00b7 Measuring change  (two families)",
                      "3 \u00b7 The steering scaling factor  \u03b1"]):
    txt(s,0.95,4.35+i*0.62,11,0.45,[{"runs":[{"t":t,"sz":17,"b":True,"font":BODY,"c":ICE}]}])
notes(s,
 "Appendix \u2014 the method behind the results. Three short slides: how every direction (axis) is built, how we measure "
 "change, and how the steering knob works. No new results here.")

# ==================================================== M1 \u00b7 THE AXES
s=slide(WHITE)
pill(s,0.7,0.45,"METHOD 1 \u00b7 THE AXES",fill=BLUE,sz=11)
txt(s,0.68,0.9,12.2,0.9,[{"runs":[{"t":"Finding the axes \u2014 one recipe: diff-in-means","sz":28,"b":True,"font":HEAD,"c":NAVY}]}])
flow=["Two contrastive sets\n(condition A vs B)","Last-token residual\nat every layer","normalize( mean_A \u2212 mean_B )\n\u2192 one direction per layer"]
fx=[0.7,4.75,8.8]
for i,(x,t) in enumerate(zip(fx,flow)):
    card(s,x,1.95,3.7,1.15,fill=CARD)
    txt(s,x+0.2,1.95,3.3,1.15,[{"runs":[{"t":t,"sz":13.5,"b":True,"font":BODY,"c":INK}],"line":1.08,"align":PP_ALIGN.CENTER}],anchor=MSO_ANCHOR.MIDDLE)
    if i<2: txt(s,x+3.7,1.95,0.35,1.15,[{"runs":[{"t":"\u2192","sz":22,"b":True,"font":BODY,"c":BLUE}]}],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
txt(s,0.7,3.45,12,0.4,[{"runs":[{"t":"The same recipe builds all four directions:","sz":15,"b":True,"font":BODY,"c":MUTED}]}])
axes=[("Image valence","distressing  vs  positive  images  (same describe prompt)",BLUE),
      ("Text valence","negative  vs  positive  sentences",GREEN),
      ("Text emotion","emotion stories  vs  neutral  (desperation, fear, \u2026)",ORANGE),
      ("Refusal  r","harmful  vs  harmless  prompts",NAVY)]
for i,(nm,desc,col) in enumerate(axes):
    yy=3.95+i*0.52
    pill(s,0.9,yy,nm,fill=col,sz=11,w=2.55)
    txt(s,3.65,yy+0.02,9.2,0.4,[{"runs":[{"t":desc,"sz":13.5,"font":BODY,"c":INK}]}])
txt(s,0.7,6.35,12.2,0.6,[{"runs":[{"t":"OASIS valence tertiles for images \u00b7 residual taken at the last token, kept per layer \u00b7 axis reproducibility checked by split-half cosine (a_stability).","sz":11,"font":BODY,"c":MUTED}],"line":1.05}])
notes(s,
 "Plain terms: every internal 'direction' is built the same way \u2014 take two opposite sets, run them through the model, and "
 "subtract their average internal states. The result is one direction per layer.\n\n"
 "Datasets, per direction: IMAGE valence \u2014 OASIS (most-distressing vs most-positive images, same describe prompt); TEXT "
 "valence and TEXT emotion \u2014 Claude-generated sentences (emotion stories vs neutral); REFUSAL \u2014 AdvBench (harmful) vs "
 "Alpaca (harmless). We keep the direction at every layer and check it's reproducible via split-half cosine.")

# ==================================================== M2 \u00b7 MEASURING CHANGE
s=slide(WHITE)
pill(s,0.7,0.45,"METHOD 2 \u00b7 MEASURING CHANGE",fill=BLUE,sz=11)
txt(s,0.68,0.9,12.2,0.9,[{"runs":[{"t":"Measuring change \u2014 two families","sz":28,"b":True,"font":HEAD,"c":NAVY}]}])
# left card A
card(s,0.7,1.95,5.85,3.7,fill=CARD)
pill(s,1.0,2.2,"A \u00b7 REPRESENTATIONAL (geometry)",fill=GREEN,sz=11)
txt(s,1.0,2.8,5.3,2.7,[
  {"runs":[{"t":"Alignment:","b":True,"sz":14,"c":INK},{"t":" per-layer cosine between two directions (e.g. text-desperation vs image-valence = 0.37).","sz":14,"c":INK}],"line":1.1,"space_after":8},
  {"runs":[{"t":"Evocation:","b":True,"sz":14,"c":INK},{"t":" projection (dot product) of an activation onto a direction \u2014 distress vs positive image.","sz":14,"c":INK}],"line":1.1,"space_after":10},
  {"runs":[{"t":"Cosine appears ONLY here.","b":True,"sz":13.5,"c":GREEN}]}])
# right card B
card(s,6.75,1.95,5.85,3.7,fill=CARD)
pill(s,7.05,2.2,"B \u00b7 BEHAVIORAL (first-token logit)",fill=BLUE,sz=11)
txt(s,7.05,2.8,5.3,2.7,[
  {"runs":[{"t":"score = logsumexp(A) \u2212 logsumexp(B)","b":True,"sz":13.5,"font":MONO,"c":NAVY}],"space_after":6},
  {"runs":[{"t":"log-odds of answer-family A over B at the first token \u2014 generation-free.","sz":14,"c":INK}],"line":1.1,"space_after":8},
  {"runs":[{"t":"Refusal","b":True,"sz":14,"c":INK},{"t":" = refuse-vs-comply tokens; rate = fraction > 0.","sz":14,"c":INK}],"line":1.1,"space_after":8},
  {"runs":[{"t":"\u0394 measured under image/text swap, steering (\u00b1\u03b1), and controls.","sz":14,"c":INK}],"line":1.1}])
callout(s,0.7,5.85,11.9,1.0,
  "Cosine is only for aligning two vectors. All behavior \u2014 including refusal \u2014 is a difference of first-token log-probabilities. No text generated, no judge.",
  fill=NAVY,tc=WHITE,sz=13.5,font=BODY)
notes(s,
 "Plain terms: two kinds of measurement. (A) GEOMETRY \u2014 how aligned two internal directions are (cosine) \u2014 used only for "
 "the 'shared representation' result. (B) BEHAVIOR \u2014 how much the model leans toward answer A vs B on its very next word "
 "(log-odds), never generating text and never using a judge. Refusal is the same idea (refuse-word odds vs comply-word "
 "odds). No dataset beyond the prompts already cited.")

# ==================================================== M3 \u00b7 SCALING FACTOR
s=slide(WHITE)
pill(s,0.7,0.45,"METHOD 3 \u00b7 STEERING",fill=BLUE,sz=11)
txt(s,0.68,0.9,12.2,0.9,[{"runs":[{"t":"The steering scaling factor  \u03b1","sz":28,"b":True,"font":HEAD,"c":NAVY}]}])
callout(s,0.7,1.9,11.9,0.95,
  "\u0394  =  \u03b1  \u00d7  ||resid||  \u00d7  \u00e2        added to resid_post at every layer",
  fill=CARD,tc=NAVY,sz=17,font=MONO)
# left: why norm-scaled
card(s,0.7,3.1,5.85,2.9,fill=CARD)
pill(s,1.0,3.35,"WHY SCALE TO ||resid|| ?",fill=BLUE,sz=11)
txt(s,1.0,3.95,5.3,2.0,[
  {"runs":[{"t":"The residual norm grows hugely with depth, so \u03b1 is a ","sz":14,"c":INK},{"t":"fraction of the local residual scale","b":True,"sz":14,"c":INK},{"t":", not an absolute size.","sz":14,"c":INK}],"line":1.12,"space_after":8},
  {"runs":[{"t":"\u03b1 = 0.01  \u2192  add 1% of the local residual magnitude along \u00e2.","b":True,"sz":14,"c":NAVY}],"line":1.12,"space_after":8},
  {"runs":[{"t":"Dimensionless \u2192 comparable across layers and models.","sz":13.5,"c":MUTED}],"line":1.1}])
# right: what alpha does
card(s,6.75,3.1,5.85,2.9,fill=CARD)
pill(s,7.05,3.35,"WHAT \u03b1 DOES (dose-response)",fill=ORANGE,sz=11)
txt(s,7.05,3.95,5.3,2.0,[
  {"runs":[{"t":"small (\u22640.006):","b":True,"sz":14,"c":INK},{"t":" gentle tilt, output intact.","sz":14,"c":INK}],"line":1.1,"space_after":6},
  {"runs":[{"t":"0.008\u20130.012:","b":True,"sz":14,"c":INK},{"t":" flips refusal / behavior, still coherent.","sz":14,"c":INK}],"line":1.1,"space_after":6},
  {"runs":[{"t":">0.016:","b":True,"sz":14,"c":VERM},{"t":" swamps computation \u2192 incoherent (gated out).","sz":14,"c":INK}],"line":1.1}])
txt(s,0.7,6.25,12.2,0.6,[{"runs":[{"t":"||resid|| includes the massive-activation dims \u2014 which is why the causal effect surviving their removal (Fig 2) rules out a scaling artifact.","sz":11,"font":BODY,"c":MUTED}],"line":1.05}])
notes(s,
 "Plain terms: to 'steer,' we add a small push along a direction at every layer, sized as a percentage of the model's own "
 "internal magnitude \u2014 so alpha is a dimensionless dial. Small pushes tilt gently; medium ones flip behavior while keeping "
 "the text coherent; too-large ones break the output (we discard those). No dataset \u2014 this is the steering mechanism itself.")

prs.save(OUT)
print("saved:",OUT,"| slides:",len(prs.slides._sldIdLst))
